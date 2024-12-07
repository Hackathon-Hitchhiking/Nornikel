from typing import Generator, Dict, Any, List
from io import BytesIO
import base64

from docx import Document
from docx.text.paragraph import Paragraph
from docx.table import Table
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

import fitz


class DocxProcessor:
    """
    Класс для обработки .docx файлов, включая извлечение текста и изображений.
    """

    def __init__(self, chunk_size: int = 100, overlap: int = 25):
        """
        Инициализация процессора.
        :param chunk_size: Максимальное количество слов в одном чанке.
        :param overlap: Количество слов, которые перекрываются между чанками.
        """
        self.chunk_size = chunk_size
        self.overlap = overlap

    def process(self, file_bytes: bytes) -> Generator[Dict[str, Any], None, None]:
        """
        Обрабатывает .docx файл и возвращает чанки текста с метаданными.
        :param file_bytes: Содержимое .docx файла в байтах.
        :return: Генератор чанков текста с метаданными.
        """
        document = self._load_document(file_bytes)
        tokens, total_words = self._extract_tokens_with_positions(document)
        yield from self._generate_chunks_with_metadata(tokens, total_words)

    def _load_document(self, file_bytes: bytes) -> Document:
        """
        Загружает .docx документ из байтов.
        :param file_bytes: Содержимое .docx файла в байтах.
        :return: Объект документа.
        """
        return Document(BytesIO(file_bytes))

    def iter_block_items(self, parent):
        """
        Возвращает параграфы и таблицы в порядке их следования в документе.
        """
        parent_elm = parent.element.body
        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)

    def run_has_image(self, run):
        """
        Проверяет, содержит ли run изображение.
        """
        r_element = run._r
        drawing_elements = r_element.findall(
            ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing"
        )
        return bool(drawing_elements)

    def extract_image_from_run(self, run):
        """
        Извлекает изображение из run.
        """
        r_element = run._r
        drawing_elements = r_element.findall(
            ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing"
        )
        for drawing in drawing_elements:
            blip_elements = drawing.findall(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
            )
            if blip_elements:
                embed = blip_elements[0].get(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                )
                image_part = run.part.related_parts[embed]
                image_bytes = image_part.blob
                return image_bytes
        return None

    def _extract_tokens_with_positions(
        self, document: Document
    ) -> (List[Dict[str, Any]], int):
        """
        Извлекает текст и изображения из документа в порядке их следования, с позициями.
        :param document: Объект документа.
        :return: Список токенов (слов и изображений) с позициями, и общее количество слов.
        """
        tokens = []
        current_word_position = 0

        for block in self.iter_block_items(document):
            if isinstance(block, Paragraph):
                paragraph = block
                for run in paragraph.runs:
                    if self.run_has_image(run):
                        image_bytes = self.extract_image_from_run(run)
                        if image_bytes:
                            tokens.append(
                                {
                                    "type": "image",
                                    "image_bytes": base64.b64encode(image_bytes).decode(
                                        "utf-8"
                                    ),
                                    "position": current_word_position,
                                }
                            )
                    if run.text.strip():
                        words = run.text.strip().split()
                        for word in words:
                            tokens.append(
                                {
                                    "type": "word",
                                    "content": word,
                                    "position": current_word_position,
                                }
                            )
                            current_word_position += 1
            elif isinstance(block, Table):
                pass

        total_words = current_word_position
        return tokens, total_words

    def _generate_chunks_with_metadata(
        self,
        tokens: List[Dict[str, Any]],
    ) -> Generator[Dict[str, Any], None, None]:
        """
        Создает чанки текста с метаданными.
        :param tokens: Список токенов (слов и изображений) с позициями.
        :param total_words: Общее количество слов в документе
        :return: Генератор чанков текста с метаданными.
        """
        idx = 0
        total_tokens = len(tokens)
        chunk_id = 0

        while idx < total_tokens:
            chunk_tokens = tokens[idx : idx + self.chunk_size]
            chunk_text = " ".join(
                token["content"] for token in chunk_tokens if token["type"] == "word"
            )
            chunk_images = [token for token in chunk_tokens if token["type"] == "image"]
            start_word = chunk_tokens[0]["position"]
            end_word = chunk_tokens[-1]["position"] + 1

            chunk_id += 1
            yield {
                "chunk": chunk_text,
                "metadata": {
                    "start_word": start_word,
                    "end_word": end_word,
                    "images": chunk_images,
                },
            }
            idx += self.chunk_size - self.overlap


class PdfProcessor:
    """
    Класс для обработки PDF файлов, включая извлечение текста и изображений.
    """

    def __init__(self, chunk_size: int = 100, overlap: int = 25):
        """
        Инициализация процессора.
        :param chunk_size: Максимальное количество слов в одном чанке.
        :param overlap: Количество слов, которые перекрываются между чанками.
        """
        self.chunk_size = chunk_size
        self.overlap = overlap

    def process(self, file_bytes: bytes) -> Generator[Dict[str, Any], None, None]:
        """
        Обрабатывает PDF файл и возвращает чанки текста с метаданными.
        :param file_bytes: Содержимое PDF файла в байтах.
        :return: Генератор чанков текста с метаданными.
        """
        tokens = self._extract_tokens_with_positions(file_bytes)
        yield from self._generate_chunks_with_metadata(tokens)

    def _extract_tokens_with_positions(self, file_bytes: bytes) -> List[Dict[str, Any]]:
        """
        Извлекает текст и изображения из PDF в порядке их следования, с позициями.
        :param file_bytes: Содержимое PDF файла в байтах.
        :return: Список токенов (слов и изображений) с позициями.
        """
        tokens = []
        current_word_position = 0

        # Открываем PDF из байтов
        pdf_document = fitz.open(stream=file_bytes, filetype="pdf")

        for page_number in range(len(pdf_document)):
            page = pdf_document[page_number]
            text = page.get_text()
            if text:
                words = text.strip().split()
                for word in words:
                    tokens.append(
                        {
                            "type": "word",
                            "content": word,
                            "position": current_word_position,
                            "page_number": page_number + 1,
                        }
                    )
                    current_word_position += 1

            # Извлечение изображений
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                xref = img[0]
                try:
                    base_image = pdf_document.extract_image(xref)
                    image_bytes = base_image["image"]
                    tokens.append(
                        {
                            "type": "image",
                            "image_bytes": base64.b64encode(image_bytes).decode(
                                "utf-8"
                            ),
                            "position": current_word_position,
                            "page_number": page_number + 1,
                        }
                    )
                except Exception as e:
                    print(
                        f"Ошибка при извлечении изображения на странице {page_number + 1}: {e}"
                    )

        pdf_document.close()
        return tokens

    def _generate_chunks_with_metadata(
        self,
        tokens: List[Dict[str, Any]],
    ) -> Generator[Dict[str, Any], None, None]:
        """
        Создает чанки текста с метаданными.
        :param tokens: Список токенов (слов и изображений) с позициями.
        :return: Генератор чанков текста с метаданными.
        """
        idx = 0
        total_tokens = len(tokens)
        chunk_id = 0

        while idx < total_tokens:
            chunk_tokens = tokens[idx : idx + self.chunk_size]
            chunk_text = " ".join(
                token["content"] for token in chunk_tokens if token["type"] == "word"
            )
            chunk_images = [token for token in chunk_tokens if token["type"] == "image"]
            start_word = chunk_tokens[0]["position"]
            end_word = chunk_tokens[-1]["position"] + 1

            chunk_id += 1
            yield {
                "chunk": chunk_text,
                "metadata": {
                    "start_word": start_word,
                    "end_word": end_word,
                    "images": chunk_images,
                },
            }
            idx += self.chunk_size - self.overlap


class PptxProcessor:
    """
    Класс для обработки .pptx файлов, включая извлечение текста и изображений.
    """

    def __init__(self, chunk_size: int = 100, overlap: int = 25):
        """
        Инициализация процессора.
        :param chunk_size: Максимальное количество слов в одном чанке.
        :param overlap: Количество слов, которые перекрываются между чанками.
        """
        self.chunk_size = chunk_size
        self.overlap = overlap

    def process(self, file_bytes: bytes) -> Generator[Dict[str, Any], None, None]:
        """
        Обрабатывает .pptx файл и возвращает чанки текста с метаданными.
        :param file_bytes: Содержимое .pptx файла в байтах.
        :return: Генератор чанков текста с метаданными.
        """
        presentation = self._load_presentation(file_bytes)
        tokens, total_words = self._extract_tokens_with_positions(presentation)
        yield from self._generate_chunks_with_metadata(tokens, total_words)

    def _load_presentation(self, file_bytes: bytes) -> Presentation:
        """
        Загружает .pptx документ из байтов.
        :param file_bytes: Содержимое .pptx файла в байтах.
        :return: Объект презентации.
        """
        return Presentation(BytesIO(file_bytes))

    def _extract_tokens_with_positions(
        self, presentation: Presentation
    ) -> (List[Dict[str, Any]], int):
        """
        Извлекает текст и изображения из презентации в порядке их следования, с позициями.
        :param presentation: Объект презентации.
        :return: Список токенов (слов и изображений) с позициями, и общее количество слов.
        """
        tokens = []
        current_word_position = 0

        for slide_number, slide in enumerate(presentation.slides, start=1):
            # Извлечение текста
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            words = run.text.strip().split()
                            for word in words:
                                tokens.append(
                                    {
                                        "type": "word",
                                        "content": word,
                                        "position": current_word_position,
                                        "slide_number": slide_number,
                                    }
                                )
                                current_word_position += 1

                # Извлечение изображений
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_bytes = image.blob
                    tokens.append(
                        {
                            "type": "image",
                            "image_bytes": base64.b64encode(image_bytes).decode(
                                "utf-8"
                            ),
                            "position": current_word_position,
                            "slide_number": slide_number,
                        }
                    )

        total_words = current_word_position
        return tokens, total_words

    def _generate_chunks_with_metadata(
        self, tokens: List[Dict[str, Any]], total_words: int
    ) -> Generator[Dict[str, Any], None, None]:
        """
        Создает чанки текста с метаданными.
        :param tokens: Список токенов (слов и изображений) с позициями.
        :param total_words: Общее количество слов в презентации.
        :return: Генератор чанков текста с метаданными.
        """
        idx = 0
        total_tokens = len(tokens)
        chunk_id = 0

        while idx < total_tokens:
            chunk_tokens = tokens[idx : idx + self.chunk_size]
            chunk_text = " ".join(
                token["content"] for token in chunk_tokens if token["type"] == "word"
            )
            chunk_images = [token for token in chunk_tokens if token["type"] == "image"]
            start_word = chunk_tokens[0]["position"]
            end_word = chunk_tokens[-1]["position"] + 1

            chunk_id += 1
            yield {
                "chunk": chunk_text,
                "metadata": {
                    "start_word": start_word,
                    "end_word": end_word,
                    "images": chunk_images,
                    "slide_numbers": list(
                        set(token["slide_number"] for token in chunk_tokens)
                    ),
                },
            }
            idx += self.chunk_size - self.overlap
